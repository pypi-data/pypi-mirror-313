import os
import sys

import spacy
from pptx import Presentation
from PyPDF2 import PdfReader


class PDFExtractor:
    """
    PDFExtractor class to extract text from PDF files
    """

    def __init__(self):
        pass

    def extract(self, file_path):
        reader = PdfReader(file_path)
        all_text = " ".join([page.extract_text() for page in reader.pages])
        all_text = all_text.replace("\n", " ")
        return all_text


class VideoExtractor:
    """
    VideoExtractor class to extract text from video files
    """

    def __init__(self):
        pass


class PPTExtractor:
    """
    PPTExtractor class to extract text from PowerPoint files
    """

    def __init__(self):
        pass

    def extract(self, file_path):
        text = []
        prs = Presentation(file_path)

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return " ".join(text)


def get_text_extractor(filename):
    """
    Get the appropriate text extractor based on the file extension
    Args:
        filename: str: name of the file
    Returns:
        extractor: object: text extractor object
    """

    if filename.lower().endswith(".pdf"):
        return PDFExtractor()
    elif filename.lower().endswith(".pptx"):
        return PPTExtractor()
    return None


def extract_text_from_files(filelist):
    """
    Extract text from a list of files
    Args:
        filelist: list: list of file paths
    Returns:
        text: list: list of extracted text
    """

    text = []
    for f in filelist:
        extractor = get_text_extractor(f)
        if extractor is None:
            text.append("")
        else:
            text.append(extractor.extract(f))
    return text


def extract_text_from_folder(folderpath):
    """
    Extract text from all files in a folder
    Args:
        folderpath: str: path to the folder
    Returns:
        text: list: list of extracted text
        files: list: list of file paths
    """

    files = os.listdir(folderpath)
    files = [
        os.path.join(folderpath, f)
        for f in files
        if os.path.isfile(os.path.join(folderpath, f))
    ]

    texts = extract_text_from_files(files)

    return (texts, files)


def tokenize_sentences(text):
    """
    Tokenize text into sentences
    Args:
        text: str or list: text to tokenize
    Returns:
        sentences: list: list of tokenized sentences
    """

    sentences = []
    nlp = spacy.load("en_core_web_sm")
    if type(text) is list:
        for t in text:
            doc = nlp(t)
            sentences.append([str(s) for s in doc.sents])
    elif type(text) is str:
        doc = nlp(t)
        sentences.append([str(s) for s in doc.sents])

    return sentences


if __name__ == "__main__":
    args = sys.argv
    fp = args[1]

    text, files = extract_text_from_folder(fp)
    print(len(text), len(files))
    print(text)
    doc_sentences = tokenize_sentences(text)
    print(f"length of tokenized sentences {len(doc_sentences)}")
    for sentences, file in zip(doc_sentences, files, strict=False):
        print(file)
        for sentence in sentences:
            print("* ", sentence)
