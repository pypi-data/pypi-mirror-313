"""Power point interface."""

from bacore.domain.measurements import Time
from dataclasses import dataclass
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.shapes.picture import Picture
from pptx.slide import Slide
from pptx.text.text import TextFrame
from pptx.util import Inches, Length, Pt
from sqlmodel import Field, SQLModel
from typing import ClassVar

TODAY = Time().today_s


@dataclass
class PowerPoint:
    """PowerPoint class encapsulating the pptx.Presentation object.

    Attributes:
        prs (Presentation): The PowerPoint presentation object.
        background_layer (ClassVar[int]): Layer index for background elements.
        widescreen_width (ClassVar[Inches]): Width for widescreen presentations.
        widescreen_height (ClassVar[Inches]): Height for widescreen presentations.
    """

    prs = Presentation()

    background_layer: ClassVar[int] = 2
    widescreen_width: ClassVar[Inches] = Inches(13.33)
    widescreen_height: ClassVar[Inches] = Inches(7.5)

    def add_slide(self, layout_index: int, title_text: str | None = None) -> Slide:
        """Add a PowerPoint slide with an optional title text.

        Arguments:
            layout_index (int): Index of the slide layout to use.
            title_text (Optional[str]): Optional title text for the slide.

        Returns:
            Slide: The newly added slide.

        Raises:
            ValueError: If the layout_index is out of range.
        """
        if layout_index < 0 or layout_index >= len(self.prs.slide_layouts):
            raise ValueError(f"Layout index '{layout_index}' out of range.")

        slide_layout = self.prs.slide_layouts[layout_index]
        slide = self.prs.slides.add_slide(slide_layout)

        if title_text and slide.shapes.title:
            title = slide.shapes.title
            title.text = title_text
        return slide

    @staticmethod
    def add_background_image(
        slide: Slide,
        image_file: str,
        left: Length = 0,
        top: Length = 0,
        width: Length | None = None,
        height: Length | None = None,
        move_to_background: bool = True,
    ) -> Picture:
        """Add a background image to a slide.

        Arguments:
            slide (Slide): The slide to which the image will be added.
            image_file (str): Path to the image file.
            left (Inches, optional): The left position of the image. Defaults to Inches(0).
            top (Inches, optional): The top position of the image. Defaults to Inches(0).
            width (Optional[Inches], optional): The width of the image. Defaults to None.
            height (Optional[Inches], optional): The height of the image. Defaults to None.
            move_to_background (bool, optional): If True, moves the image to the background. Defaults to True.

        Returns:
            Picture: The added picture object.
        """
        background_img = slide.shapes.add_picture(image_file, left, top, width, height)
        if move_to_background:
            slide.shapes._spTree.remove(background_img._element)
            slide.shapes._spTree.insert(2, background_img._element)
        return background_img

    @classmethod
    def default_templates(cls, save_path: str = "default_templates.pptx", widescreen: bool = True) -> None:
        """Create a PowerPoint presentation using default templates.

        Arguments:
            save_path (str, optional): Path to save the default templates. Defaults to "default_templates.pptx".
            widescreen (bool, optional): Whether to set the presentation to widescreen. Defaults to True.
        """
        ppt = cls()

        if widescreen:
            ppt.prs.slide_width = cls.widescreen_width
            ppt.prs.slide_height = cls.widescreen_height

        for layout_index in range(len(ppt.prs.slide_layouts)):
            ppt.add_slide(layout_index)

        ppt.prs.save(save_path)
        print(f"Default templates saved to {save_path}")


@dataclass
class Placeholder:
    """Represents a placeholder item inside a slide.

    Attributes:
        slide (Slide): The slide containing the placeholder.
        number (int): The placeholder index.
    """

    slide: Slide
    number: int


@dataclass
class Bullet:
    """Represents a text bullet with optional bullet level and URL.

    Attributes:
        text (str): The bullet text.
        bullet_level (Optional[int]): The indentation level of the bullet.
        font_size (int): The font size of the bullet text.
        url (Optional[str]): An optional hyperlink for the bullet.
    """

    text: str
    bullet_level: int | None = None
    font_size: int = 14
    url: str | None = None


class TitleSlide(SQLModel, table=True):
    """Represents the title slide of a presentation."""

    __tablename__ = "title_slide"

    id: int | None = Field(default=None, primary_key=True)
    title: str = Field(index=True)
    sub_title: str | None = Field(default=None)
    background_image: str | None = Field(default=None)
    date: str | None = Field(default=None)
    logo: str | None = Field(default=None)

    def create(self, ppt: PowerPoint) -> Slide:
        """Create a title slide in the PowerPoint presentation.

        Arguments:
            ppt (PowerPoint): The PowerPoint presentation object.

        Returns:
            Slide: The created title slide.
        """
        slide = ppt.add_slide(layout_index=0, title_text=self.title)
        if self.background_image:
            PowerPoint.add_background_image(
                slide=slide,
                image_file=self.background_image,
                width=PowerPoint.widescreen_width,
                height=Inches(7),
            )
        if self.sub_title:
            subtitle = slide.shapes.placeholders[1]
            subtitle.text = self.sub_title
            subtitle.text_frame.paragraphs[0].font.italic = True
            subtitle.text_frame.paragraphs[0].font.color.rbg = RGBColor(255, 0, 0)
        if self.date:
            date_box = slide.shapes.add_textbox(left=Inches(0.5), top=Inches(7.1), width=Inches(1), height=Inches(0.3))
            date_tf = date_box.text_frame
            date_tf.text = self.date
            date_p = date_tf.paragraphs[0]
            date_run = date_p.runs[0]
            date_run.font.size = Pt(12)
            date_run.font.italic = True
            date_p.alignment = PP_ALIGN.CENTER
            date_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        if self.logo:
            slide.shapes.add_picture(
                self.logo,
                left=Inches(11),
                top=Inches(7),
                width=Inches(1.5),
            )
        return slide


def create_bullets(container: Placeholder | TextFrame, bullets: list[Bullet]):
    """Add bullets to a placeholder or text frame.

    Arguments:
        container (Union[Placeholder, TextFrame]): The container to add bullets to.
        bullets (List[Bullet]): A list of Bullet instances to add.
    """
    if isinstance(container, Placeholder):
        placeholder = container.slide.shapes.placeholders[container.number]
        text_frame = placeholder.text_frame
    else:
        text_frame = container

    for bullet in bullets:
        paragraph = text_frame.add_paragraph()
        if bullet.url:
            paragraph = paragraph.add_run()
            paragraph.text = bullet.text
            paragraph.hyperlink.address = bullet.url
        else:
            paragraph.text = bullet.text
        paragraph.font.size = Pt(bullet.font_size)
        if bullet.bullet_level is not None:
            paragraph.level = bullet.bullet_level
